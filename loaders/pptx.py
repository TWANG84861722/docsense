"""PowerPoint (.pptx) parsing."""
from .common import _rows_to_md
from .office_images import image_elements


def _pptx_table_md(table):
    # Takes: a pptx table object   Returns: a Markdown table string
    # Pull each row/cell's .text into a 2-D list → hand it to the shared _rows_to_md
    return _rows_to_md([[c.text for c in row.cells] for row in table.rows])


def parse_pptx(path):
    """Parse PPT (.pptx) → elements. Each slide's text → text (section='Slide N', includes notes);
    tables → table."""
    from pptx import Presentation        # lazy import: importing this module works even without python-pptx installed

    prs = Presentation(str(path))        # Takes: a file path   Returns: the whole presentation object
    elements = []
    # enumerate(prs.slides, start=1)  Returns: each (page number i, slide object), i starting at 1
    for i, slide in enumerate(prs.slides, start=1):
        section = f"Slide {i}"           # this slide's section is simply "Slide 1/2/3..."
        texts = []                       # collect all text-box text on this slide
        for shape in slide.shapes:       # iterate every "shape" on the slide (text box / table / image ...)
            if shape.has_table:                       # this shape is a table
                md = _pptx_table_md(shape.table)      # → a Markdown table
                if md:
                    elements.append({"page": i, "section": section, "type": "table", "text": md})
            elif shape.has_text_frame:                # this shape is a text box
                t = shape.text_frame.text.strip()     # grab its text
                if t:
                    texts.append(t)
        if slide.has_notes_slide:                     # if the slide has speaker notes, collect them too
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                texts.append(f"[Notes] {note}")
        if texts:                                     # merge this slide's collected text into one text element
            elements.append({"page": i, "section": section, "type": "text", "text": "\n".join(texts)})
    elements.extend(image_elements(path, "Slide image"))   # images embedded in slides → VL understanding → figure
    return elements
